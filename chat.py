from fastapi import APIRouter, UploadFile, Form, FastAPI
from pymongo import MongoClient
from bson import ObjectId
from datetime import datetime
from dotenv import load_dotenv
import os
import docx, io, asyncio
from pdf2image import convert_from_bytes
from pptx import Presentation
import openpyxl
import xlrd
import pytesseract
from PIL import Image

# ==========================================================
# Setup
# ==========================================================
load_dotenv()
MONGO_URI = os.getenv("MONGO_URI")

if not MONGO_URI:
    raise ValueError("❌ MONGO_URI not found in environment variables")

client = MongoClient(MONGO_URI)
db = client["chatpdf"]
chats_collection = db["chats"]

router = APIRouter(prefix="/chat", tags=["chat"])

# ==========================================================
# OCR Helper
# ==========================================================
def ocr_image(image: Image.Image) -> str:
    """Perform OCR using Tesseract on a PIL image."""
    text = pytesseract.image_to_string(image, lang="eng")
    return text.strip()

# ==========================================================
# File Text Extraction (same logic)
# ==========================================================
async def extract_text_from_file(file: UploadFile) -> tuple[str, int]:
    filename = file.filename.lower()

    if filename.endswith(".docx"):
        d = docx.Document(file.file)
        text = "\n".join([p.text for p in d.paragraphs])
        return text, 1

    elif filename.endswith(".txt"):
        try:
            text = file.file.read().decode("utf-8")
        except UnicodeDecodeError:
            file.file.seek(0)
            text = file.file.read().decode("latin-1")
        return text, 1

    elif filename.endswith(".pdf"):
        file.file.seek(0)
        pdf_bytes = file.file.read()
        pages = convert_from_bytes(pdf_bytes, dpi=100)
        text = ""
        for idx, page in enumerate(pages, start=1):
            text += f"\n--- Page {idx} ---\n{ocr_image(page)}\n"
            del page
        return text.strip(), len(pages)

    elif filename.endswith((".png", ".jpg", ".jpeg")):
        file.file.seek(0)
        image = Image.open(file.file)
        text = ocr_image(image)
        return text, 1

    elif filename.endswith(".pptx"):
        file.file.seek(0)
        prs = Presentation(file.file)
        text = ""
        for idx, slide in enumerate(prs.slides, start=1):
            text += f"\n--- Slide {idx} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        text += " | ".join([cell.text.strip() for cell in row.cells]) + "\n"
        return text.strip(), len(prs.slides)

    elif filename.endswith(".xlsx"):
        file.file.seek(0)
        wb = openpyxl.load_workbook(file.file, read_only=True)
        text = ""
        for sheet in wb.sheetnames:
            text += f"\n--- Sheet: {sheet} ---\n"
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
        return text.strip(), len(wb.sheetnames)

    elif filename.endswith(".xls"):
        file.file.seek(0)
        wb = xlrd.open_workbook(file_contents=file.file.read())
        text = ""
        for sheet in wb.sheets():
            text += f"\n--- Sheet: {sheet.name} ---\n"
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                text += " | ".join([str(cell.value) for cell in row]) + "\n"
        return text.strip(), wb.nsheets

    else:
        return "Unsupported file format.", 0

# ==========================================================
# Routes
# ==========================================================
@router.post("/create")
async def create_chat(user_email: str = Form(...), file: UploadFile = Form(...)):
    extracted_text, page_count = await extract_text_from_file(file)

    chat_doc = {
        "user_email": user_email,
        "file_name": file.filename,
        "context": extracted_text,
        "messages": [],
        "created_at": datetime.utcnow().isoformat()
    }

    result = chats_collection.insert_one(chat_doc)
    return {
        "status": "success",
        "chat_id": str(result.inserted_id),
        "page_count": page_count,
        "text_length": len(extracted_text),
        "pages_snippet": extracted_text[:1000],
    }

@router.get("/getall")
def get_all_chats(user_email: str):
    chats = chats_collection.find({"user_email": user_email}, {"file_name": 1, "created_at": 1})
    result = [{"chat_id": str(c["_id"]), "file_name": c["file_name"], "created_at": c.get("created_at")} for c in chats]
    return {"status": "success", "chats": result}

@router.get("/history/{chat_id}")
def get_chat_history(chat_id: str):
    chat = chats_collection.find_one({"_id": ObjectId(chat_id)}, {"messages": 1, "file_name": 1})
    if not chat:
        return {"status": "error", "message": "Chat not found ❌"}
    return {"status": "success", "file_name": chat["file_name"], "messages": chat.get("messages", [])}

@router.delete("/delete/{chat_id}")
def delete_chat(chat_id: str, user_email: str):
    chat = chats_collection.find_one({"_id": ObjectId(chat_id)})
    if not chat:
        return {"status": "error", "message": "Chat not found ❌"}
    if chat["user_email"] != user_email:
        return {"status": "error", "message": "Not authorized ❌"}
    chats_collection.delete_one({"_id": ObjectId(chat_id)})
    return {"status": "success", "message": f"Chat {chat_id} deleted ✅"}
